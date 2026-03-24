"""
Document Text Extraction
Extracts readable text from uploaded files (PDF, DOCX, TXT, PPTX).
"""

import os


def extract_text(file_path: str) -> str:
    """
    Extract text content from a file.  Supports PDF, DOCX, TXT, and PPTX.
    Returns empty string on failure — never raises.
    """
    if not os.path.isfile(file_path):
        return ""

    ext = file_path.rsplit(".", 1)[-1].lower() if "." in file_path else ""

    extractors = {
        "pdf": _extract_pdf,
        "txt": _extract_txt,
        "md": _extract_txt,
        "docx": _extract_docx,
        "pptx": _extract_pptx,
    }

    extractor = extractors.get(ext)
    if extractor is None:
        return ""

    try:
        return extractor(file_path)
    except Exception:
        return ""


# ─── individual extractors ───

def _extract_pdf(path: str) -> str:
    from PyPDF2 import PdfReader

    pages = []
    reader = PdfReader(path)
    for page in reader.pages:
        text = page.extract_text()
        if text:
            pages.append(text)

    if pages:
        return "\n\n".join(pages)

    # Fallback to pdfplumber for PDFs where PyPDF2 yields no text.
    try:
        import pdfplumber
    except Exception:
        return ""

    with pdfplumber.open(path) as pdf:
        for page in pdf.pages:
            text = page.extract_text()
            if text:
                pages.append(text)

    if pages:
        return "\n\n".join(pages)

    # Final fallback: OCR for scanned PDFs.
    return _extract_pdf_ocr(path)


def _extract_pdf_ocr(path: str) -> str:
    try:
        import fitz  # PyMuPDF
        import pytesseract
        from PIL import Image
    except Exception:
        return ""

    tesseract_cmd = os.getenv("TESSERACT_CMD")
    if not tesseract_cmd:
        default_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"
        if os.path.exists(default_cmd):
            tesseract_cmd = default_cmd
    if tesseract_cmd and os.path.exists(tesseract_cmd):
        pytesseract.pytesseract.tesseract_cmd = tesseract_cmd

    text_parts = []
    max_pages = int(os.getenv("OCR_MAX_PAGES", "5"))
    scale = float(os.getenv("OCR_SCALE", "1.5"))

    doc = fitz.open(path)
    for page_index, page in enumerate(doc):
        if page_index >= max_pages:
            break
        pix = page.get_pixmap(matrix=fitz.Matrix(scale, scale))
        mode = "RGB" if pix.alpha == 0 else "RGBA"
        img = Image.frombytes(mode, [pix.width, pix.height], pix.samples)
        text = pytesseract.image_to_string(img)
        if text:
            text_parts.append(text)
    doc.close()

    return "\n\n".join(text_parts)


def _extract_txt(path: str) -> str:
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()


def _extract_docx(path: str) -> str:
    import docx

    doc = docx.Document(path)
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def _extract_pptx(path: str) -> str:
    from pptx import Presentation

    prs = Presentation(path)
    text_parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_parts.append(shape.text_frame.text)
    return "\n\n".join(text_parts)
